"""Generate NadaRuns user guides (Driver + Shipper, EN + FI) as PowerPoint decks."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand palette
DARK = RGBColor(0x0C, 0x4A, 0x42)      # deep teal
GREEN = RGBColor(0x10, 0xB9, 0x81)     # primary green
INK = RGBColor(0x11, 0x18, 0x27)       # near-black text
GREY = RGBColor(0x4B, 0x55, 0x63)      # body grey
LIGHT = RGBColor(0xF3, 0xF4, 0xF6)     # light fill
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SW, SH = Inches(13.333), Inches(7.5)   # 16:9

OUT_DIR = os.path.dirname(os.path.abspath(__file__))


def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _box(slide, l, t, w, h, fill=None, line=None, radius=False):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shp_type, l, t, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    s.shadow.inherit = False
    return s


def _text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          space_after=8, line_spacing=1.05):
    """runs: list of paragraphs; each paragraph is list of (text, size, bold, color)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, bold, color) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Calibri"
    return tb


def title_slide(prs, brand, title, subtitle, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DARK)
    # green accent bar
    _box(slide, 0, Inches(3.05), SW, Inches(0.10), fill=GREEN)
    # logo badge
    badge = _box(slide, Inches(0.8), Inches(0.8), Inches(0.95), Inches(0.95),
                 fill=GREEN, radius=True)
    bt = badge.text_frame
    bt.word_wrap = False
    bt.paragraphs[0].alignment = PP_ALIGN.CENTER
    bt.vertical_anchor = MSO_ANCHOR.MIDDLE
    rr = bt.paragraphs[0].add_run(); rr.text = "\u26A1"; rr.font.size = Pt(34); rr.font.color.rgb = WHITE
    _text(slide, Inches(1.9), Inches(0.95), Inches(6), Inches(0.7),
          [[("NadaRuns", 30, True, WHITE)]])
    _text(slide, Inches(0.8), Inches(3.4), Inches(11.7), Inches(2.0),
          [[(title, 46, True, WHITE)]])
    _text(slide, Inches(0.82), Inches(4.95), Inches(11.5), Inches(1.0),
          [[(subtitle, 22, False, GREEN)]])
    _text(slide, Inches(0.82), Inches(6.7), Inches(11.5), Inches(0.5),
          [[(footer, 14, False, RGBColor(0xC8, 0xD6, 0xD2))]])
    return slide


def section_slide(prs, brand, step_label, heading, bullets, tip=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    # left accent column
    _box(slide, 0, 0, Inches(0.22), SH, fill=GREEN)
    # step chip
    chip = _box(slide, Inches(0.7), Inches(0.6), Inches(2.2), Inches(0.5),
                fill=LIGHT, radius=True)
    ct = chip.text_frame; ct.vertical_anchor = MSO_ANCHOR.MIDDLE
    ct.margin_top = 0; ct.margin_bottom = 0
    cp = ct.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run(); cr.text = step_label; cr.font.size = Pt(13); cr.font.bold = True; cr.font.color.rgb = DARK
    # heading
    _text(slide, Inches(0.7), Inches(1.25), Inches(12.0), Inches(1.0),
          [[(heading, 34, True, INK)]])
    # bullets
    paras = []
    for b in bullets:
        paras.append([("\u25B8  ", 17, True, GREEN), (b, 17, False, GREY)])
    _text(slide, Inches(0.8), Inches(2.45), Inches(11.8), Inches(4.0),
          paras, space_after=11, line_spacing=1.08)
    # tip
    if tip:
        box = _box(slide, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.78),
                   fill=RGBColor(0xE6, 0xF8, 0xF1), radius=True)
        tt = box.text_frame; tt.vertical_anchor = MSO_ANCHOR.MIDDLE
        tt.word_wrap = True
        tt.margin_left = Inches(0.2); tt.margin_right = Inches(0.2)
        tp = tt.paragraphs[0]
        r1 = tp.add_run(); r1.text = (brand["tip"] + "  "); r1.font.size = Pt(13); r1.font.bold = True; r1.font.color.rgb = DARK
        r2 = tp.add_run(); r2.text = tip; r2.font.size = Pt(13); r2.font.color.rgb = GREY
    return slide


def build(brand, title, subtitle, sections, filename):
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    title_slide(prs, brand, title, subtitle, brand["footer"])
    total = len(sections)
    for i, (heading, bullets, tip) in enumerate(sections, start=1):
        step_label = f'{brand["step"]} {i} / {total}'
        section_slide(prs, brand, step_label, heading, bullets, tip)
    path = os.path.join(OUT_DIR, filename)
    prs.save(path)
    print("saved", path)


# ======================= CONTENT =======================

EN = {"step": "STEP", "tip": "TIP:", "footer": "NadaRuns \u2014 the marketplace for empty transport capacity  \u00B7  Support: care@nadaruns.com"}
FI = {"step": "VAIHE", "tip": "VINKKI:", "footer": "NadaRuns \u2014 tyhj\u00E4n kuljetuskapasiteetin markkinapaikka  \u00B7  Tuki: care@nadaruns.com"}

# ---------- DRIVER EN ----------
driver_en = [
    ("What is NadaRuns?", [
        "NadaRuns is a B2B logistics marketplace that matches empty vehicle capacity and return trips with freight that needs to move.",
        "As a driver or carrier, you earn from kilometers you would otherwise drive empty.",
        "Works for solo drivers and for fleet/transport companies managing multiple vehicles.",
    ], "You only need a smartphone, a vehicle, and valid documents to start."),
    ("Download & Register", [
        "Install the NadaRuns Driver app and tap 'Register'.",
        "Choose your account type: Individual driver or Fleet / Company.",
        "Enter your name, contact details and (for fleets) your company name and business ID.",
        "Select your vehicle type (van, truck, semi-truck, tanker, etc.), licence class and registration plate.",
    ], "Fleet owners automatically get a company workspace to manage drivers and vehicles."),
    ("Verify Your Identity (KYC)", [
        "Upload the required documents: driving licence, vehicle registration and insurance.",
        "Our team reviews and approves your account, usually within a short time.",
        "You'll receive an email once you are approved (or if something needs fixing).",
    ], "You must be approved before you can accept jobs \u2014 this keeps the marketplace trusted."),
    ("Go Online & See Your Dashboard", [
        "Toggle 'Online' on the home screen to start receiving job requests.",
        "Your dashboard shows today's earnings, completed deliveries and acceptance rate.",
        "The live map shows your location and nearby available jobs.",
    ], "Toggle 'Offline' anytime to stop receiving new requests."),
    ("Receive & Accept Jobs", [
        "New job requests appear as a card over the map with pickup, drop-off, distance, payout and ETA.",
        "Review the details and tap Accept to take the job, or Decline to pass.",
        "Once accepted, the job is reserved for you and navigation begins.",
    ], "Higher acceptance rates and ratings help you get more offers."),
    ("Pick Up the Cargo", [
        "Follow the in-app route to the pickup location.",
        "On arrival, confirm pickup by entering the 4-digit pickup code (OTP).",
        "Take a Proof-of-Pickup photo of the loaded cargo.",
    ], "The OTP and photo protect both you and the shipper against disputes."),
    ("Deliver the Cargo", [
        "Navigate to the drop-off location shown in the app.",
        "Confirm delivery with the 4-digit delivery code (OTP).",
        "Take a Proof-of-Delivery photo to complete the job.",
    ], "After delivery you'll see an animated earnings summary for the trip."),
    ("Earnings & Wallet", [
        "The Wallet tab shows your Available and Pending balances.",
        "Earnings clear from Pending to Available after a short clearance window.",
        "Request a withdrawal (cash out) when you're ready \u2014 an invoice/receipt PDF is generated automatically.",
    ], "All payment receipts and withdrawal documents are emailed and stored for you."),
    ("Delivery History", [
        "The History tab lists every completed delivery with date, route and payout.",
        "Use it to track your lifetime earnings and performance.",
    ], None),
    ("Fleet Management (Company Owners)", [
        "Open Settings \u2192 Fleet to manage your company.",
        "Add and invite drivers, add vehicles and assign vehicles to drivers.",
        "Choose a job-acceptance mode: self-accept, owner-assign or hybrid.",
        "Track all company jobs, the company wallet and request payouts for approval.",
    ], "Invited drivers simply log in with their own email and password."),
    ("Settings, Language & Support", [
        "Update your profile, vehicle and notification preferences in Settings.",
        "Switch the app language between English and Finnish (Suomi) anytime.",
        "Find 'Help & Support' and 'Privacy & Terms' in your profile.",
        "Need help? Contact care@nadaruns.com.",
    ], "Drive safe and turn those empty kilometers into income!"),
]

# ---------- SHIPPER EN ----------
shipper_en = [
    ("What is NadaRuns?", [
        "NadaRuns is a B2B logistics marketplace for moving freight using available transport capacity and return trips.",
        "As a business (shipper), you get faster, more competitive transport by tapping into vehicles that are already on the road.",
        "Ideal for pallets, part-loads, full loads and special cargo across Finland.",
    ], "Less empty driving means lower cost for you and lower emissions overall."),
    ("Register Your Business", [
        "Install the NadaRuns app and register as a Business / Shipper.",
        "Enter your company name and contact details.",
        "Log in to reach your shipper dashboard.",
    ], None),
    ("Create a Shipment", [
        "Tap 'New' to start a shipment.",
        "Set pickup and drop-off locations on the map.",
        "Enter cargo details: weight, dimensions, pallet count and loading meters.",
        "Choose the required vehicle type and urgency.",
    ], "Accurate cargo details give you the most accurate price and the right vehicle."),
    ("Understand Your Price", [
        "Pricing follows the Finnish road-freight model based on chargeable freight weight.",
        "Chargeable weight is the greatest of actual weight, volume, pallet and loading-meter weight.",
        "The breakdown shows freight fee, distance, urgency and fuel surcharge.",
        "You get an instant estimate before you confirm.",
    ], "Expand the 'Advanced' section to fine-tune pallets and dimensions."),
    ("Choose How to Pay", [
        "Pay Now: pay securely by card (Stripe) right away.",
        "Accept Invoice: receive a Net-14 invoice (PDF) and pay later.",
        "Pay Later: confirm the shipment and settle afterwards.",
        "Saved cards enable one-tap instant payment for future orders.",
    ], "Funds are captured on delivery, so you only pay for completed transport."),
    ("Track Your Shipment", [
        "Watch live status updates: driver assigned, en route to pickup, picked up, delivering, delivered.",
        "See your assigned driver and estimated arrival on the map.",
    ], "You'll also get email updates at key milestones."),
    ("Proof of Pickup & Delivery", [
        "Each shipment shows a 'Proof of Pickup & Delivery' card.",
        "View the photo taken at pickup and the photo taken at delivery.",
        "Tap any photo to open it full screen.",
    ], "POP/POD photos stay available in your order summary for your records."),
    ("Invoices & Receipts", [
        "Payment receipts and invoices are generated automatically as PDFs.",
        "They are emailed to you and stored in the app for download anytime.",
    ], "Great for accounting \u2014 every transaction has a document."),
    ("Saved Payment Methods", [
        "Add and manage your cards under Settings \u2192 Payment Methods.",
        "Set a default card and remove old ones.",
        "Use one-tap 'Pay with saved card' for faster checkout.",
    ], None),
    ("Settings, Language & Support", [
        "Manage your business profile and preferences in Settings.",
        "Switch the app language between English and Finnish (Suomi) anytime.",
        "Find 'Help & Support' and 'Privacy & Terms' in your profile.",
        "Need help? Contact care@nadaruns.com.",
    ], "Move more freight, with less empty space on the road."),
]

# ---------- DRIVER FI ----------
driver_fi = [
    ("Mik\u00E4 on NadaRuns?", [
        "NadaRuns on B2B-logistiikan markkinapaikka, joka yhdist\u00E4\u00E4 vapaan ajoneuvokapasiteetin ja paluukuljetukset rahtiin, joka pit\u00E4\u00E4 siirt\u00E4\u00E4.",
        "Kuljettajana tai kuljetusyrityksen\u00E4 ansaitset kilometreist\u00E4, jotka muuten ajaisit tyhj\u00E4n\u00E4.",
        "Sopii sek\u00E4 yksitt\u00E4isille kuljettajille ett\u00E4 useita ajoneuvoja hallinnoiville kalustoyrityksille.",
    ], "Tarvitset vain \u00E4lypuhelimen, ajoneuvon ja voimassa olevat asiakirjat aloittaaksesi."),
    ("Lataa sovellus ja rekister\u00F6idy", [
        "Asenna NadaRuns Driver -sovellus ja paina 'Rekister\u00F6idy'.",
        "Valitse tilityyppi: Yksitt\u00E4inen kuljettaja tai Kalusto / Yritys.",
        "Sy\u00F6t\u00E4 nimesi ja yhteystietosi sek\u00E4 (kalustolle) yrityksen nimi ja Y-tunnus.",
        "Valitse ajoneuvotyyppi (paketti-, kuorma-, puoliper\u00E4vaunu, s\u00E4ili\u00F6 jne.), ajokorttiluokka ja rekisterinumero.",
    ], "Kalustoyrityksen omistaja saa automaattisesti ty\u00F6tilan kuljettajien ja ajoneuvojen hallintaan."),
    ("Vahvista henkil\u00F6llisyytesi (KYC)", [
        "Lataa vaaditut asiakirjat: ajokortti, ajoneuvon rekister\u00F6inti ja vakuutus.",
        "Tiimimme tarkistaa ja hyv\u00E4ksyy tilisi yleens\u00E4 nopeasti.",
        "Saat s\u00E4hk\u00F6postin, kun tilisi on hyv\u00E4ksytty (tai jos jotain pit\u00E4\u00E4 korjata).",
    ], "Sinun on oltava hyv\u00E4ksytty ennen kuin voit ottaa keikkoja \u2014 t\u00E4m\u00E4 pit\u00E4\u00E4 markkinapaikan luotettavana."),
    ("Mene linjoille ja n\u00E4e koontin\u00E4kym\u00E4", [
        "Kytke 'Online' aloitusn\u00E4yt\u00F6ll\u00E4 alkaaksesi vastaanottaa keikkapyynt\u00F6j\u00E4.",
        "Koontin\u00E4kym\u00E4 n\u00E4ytt\u00E4\u00E4 p\u00E4iv\u00E4n ansiot, toimitukset ja hyv\u00E4ksymisasteen.",
        "Reaaliaikainen kartta n\u00E4ytt\u00E4\u00E4 sijaintisi ja l\u00E4hell\u00E4 olevat keikat.",
    ], "Voit kytke\u00E4 'Offline' milloin tahansa lopettaaksesi uudet pyynn\u00F6t."),
    ("Vastaanota ja hyv\u00E4ksy keikkoja", [
        "Uudet keikkapyynn\u00F6t n\u00E4kyv\u00E4t kortteina kartalla: nouto, toimitus, et\u00E4isyys, palkkio ja arvioitu aika.",
        "Tarkista tiedot ja paina Hyv\u00E4ksy ottaaksesi keikan tai Hylk\u00E4\u00E4 ohittaaksesi.",
        "Hyv\u00E4ksynn\u00E4n j\u00E4lkeen keikka varataan sinulle ja navigointi alkaa.",
    ], "Korkea hyv\u00E4ksymisaste ja arviot tuovat sinulle lis\u00E4\u00E4 tarjouksia."),
    ("Nouda rahti", [
        "Seuraa sovelluksen reitti\u00E4 noutopaikkaan.",
        "Per\u00E4ill\u00E4 vahvista nouto sy\u00F6tt\u00E4m\u00E4ll\u00E4 4-numeroinen noutokoodi (OTP).",
        "Ota noutotodiste-valokuva lastatusta rahdista.",
    ], "OTP ja valokuva suojaavat sek\u00E4 sinua ett\u00E4 l\u00E4hett\u00E4j\u00E4\u00E4 erimielisyyksilt\u00E4."),
    ("Toimita rahti", [
        "Navigoi sovelluksen n\u00E4ytt\u00E4m\u00E4\u00E4n toimituspaikkaan.",
        "Vahvista toimitus 4-numeroisella toimituskoodilla (OTP).",
        "Ota toimitustodiste-valokuva viimeistell\u00E4ksesi keikan.",
    ], "Toimituksen j\u00E4lkeen n\u00E4et matkan ansioiden yhteenvedon."),
    ("Ansiot ja lompakko", [
        "Lompakko-v\u00E4lilehti n\u00E4ytt\u00E4\u00E4 K\u00E4ytett\u00E4viss\u00E4 olevan ja Odottavan saldon.",
        "Ansiot siirtyv\u00E4t Odottavasta K\u00E4ytett\u00E4viss\u00E4 olevaan lyhyen selvitysajan j\u00E4lkeen.",
        "Pyyd\u00E4 nostoa, kun olet valmis \u2014 lasku/kuittipdf luodaan automaattisesti.",
    ], "Kaikki maksukuitit ja nostoasiakirjat l\u00E4hetet\u00E4\u00E4n s\u00E4hk\u00F6postiisi ja tallennetaan."),
    ("Toimitushistoria", [
        "Historia-v\u00E4lilehti listaa jokaisen toimituksen p\u00E4iv\u00E4m\u00E4\u00E4r\u00E4n, reitin ja palkkion.",
        "Seuraa sen avulla kokonaisansioitasi ja suoritustasi.",
    ], None),
    ("Kaluston hallinta (yrityksen omistajat)", [
        "Avaa Asetukset \u2192 Kalusto hallitaksesi yrityst\u00E4si.",
        "Lis\u00E4\u00E4 ja kutsu kuljettajia, lis\u00E4\u00E4 ajoneuvoja ja osoita ajoneuvot kuljettajille.",
        "Valitse keikkojen hyv\u00E4ksymistapa: itse-hyv\u00E4ksynt\u00E4, omistajan osoitus tai hybridi.",
        "Seuraa yrityksen keikkoja ja lompakkoa sek\u00E4 pyyd\u00E4 maksatuksia hyv\u00E4ksytt\u00E4v\u00E4ksi.",
    ], "Kutsutut kuljettajat kirjautuvat sis\u00E4\u00E4n omalla s\u00E4hk\u00F6postilla ja salasanalla."),
    ("Asetukset, kieli ja tuki", [
        "P\u00E4ivit\u00E4 profiilisi, ajoneuvosi ja ilmoitusasetukset Asetuksissa.",
        "Vaihda sovelluksen kieli englannin ja suomen v\u00E4lill\u00E4 milloin tahansa.",
        "L\u00F6yd\u00E4t 'Ohje ja tuki' sek\u00E4 'Tietosuoja ja ehdot' profiilistasi.",
        "Tarvitsetko apua? Ota yhteytt\u00E4: care@nadaruns.com.",
    ], "Aja turvallisesti ja muuta tyhj\u00E4t kilometrit tuloiksi!"),
]

# ---------- SHIPPER FI ----------
shipper_fi = [
    ("Mik\u00E4 on NadaRuns?", [
        "NadaRuns on B2B-logistiikan markkinapaikka rahdin siirt\u00E4miseen vapaan kuljetuskapasiteetin ja paluukuljetusten avulla.",
        "Yrityksen\u00E4 (l\u00E4hett\u00E4j\u00E4n\u00E4) saat nopeampaa ja kilpailukykyisemp\u00E4\u00E4 kuljetusta hy\u00F6dynt\u00E4m\u00E4ll\u00E4 jo liikkeell\u00E4 olevia ajoneuvoja.",
        "Sopii lavoille, osakuormille, t\u00E4ysille kuormille ja erikoisrahdille ymp\u00E4ri Suomen.",
    ], "V\u00E4hemm\u00E4n tyhj\u00E4n\u00E4ajoa tarkoittaa pienempi\u00E4 kustannuksia ja v\u00E4hemm\u00E4n p\u00E4\u00E4st\u00F6j\u00E4."),
    ("Rekister\u00F6i yrityksesi", [
        "Asenna NadaRuns-sovellus ja rekister\u00F6idy Yritys / L\u00E4hett\u00E4j\u00E4 -tilin\u00E4.",
        "Sy\u00F6t\u00E4 yrityksesi nimi ja yhteystiedot.",
        "Kirjaudu sis\u00E4\u00E4n p\u00E4\u00E4st\u00E4ksesi l\u00E4hett\u00E4j\u00E4n koontin\u00E4kym\u00E4\u00E4n.",
    ], None),
    ("Luo kuljetus", [
        "Paina 'Uusi' aloittaaksesi kuljetuksen.",
        "Aseta nouto- ja toimituspaikat kartalle.",
        "Sy\u00F6t\u00E4 rahdin tiedot: paino, mitat, lavojen m\u00E4\u00E4r\u00E4 ja lavametrit.",
        "Valitse tarvittava ajoneuvotyyppi ja kiireellisyys.",
    ], "Tarkat rahtitiedot antavat tarkimman hinnan ja oikean ajoneuvon."),
    ("Ymm\u00E4rr\u00E4 hinta", [
        "Hinnoittelu noudattaa suomalaista maantiekuljetusmallia rahdituspainon perusteella.",
        "Rahdituspaino on suurin seuraavista: todellinen paino, tilavuus, lava- ja lavametripaino.",
        "Erittely n\u00E4ytt\u00E4\u00E4 rahtimaksun, et\u00E4isyyden, kiireellisyyden ja polttoainelis\u00E4n.",
        "Saat v\u00E4litt\u00F6m\u00E4n arvion ennen vahvistusta.",
    ], "Avaa 'Lis\u00E4asetukset' hienos\u00E4\u00E4t\u00E4\u00E4ksesi lavoja ja mittoja."),
    ("Valitse maksutapa", [
        "Maksa heti: maksa turvallisesti kortilla (Stripe) saman tien.",
        "Hyv\u00E4ksy lasku: saat Netto-14 -laskun (PDF) ja maksat my\u00F6hemmin.",
        "Maksa my\u00F6hemmin: vahvista kuljetus ja maksa j\u00E4lkik\u00E4teen.",
        "Tallennetut kortit mahdollistavat yhden napautuksen pikamaksun jatkossa.",
    ], "Varat veloitetaan vasta toimituksen yhteydess\u00E4, joten maksat vain valmiista kuljetuksesta."),
    ("Seuraa kuljetustasi", [
        "Seuraa reaaliaikaisia tilap\u00E4ivityksi\u00E4: kuljettaja osoitettu, matkalla noutoon, noudettu, toimitetaan, toimitettu.",
        "N\u00E4e osoitettu kuljettaja ja arvioitu saapuminen kartalla.",
    ], "Saat my\u00F6s s\u00E4hk\u00F6posti-ilmoituksia t\u00E4rkeiss\u00E4 vaiheissa."),
    ("Nouto- ja toimitustodiste", [
        "Jokaisessa kuljetuksessa on 'Nouto- ja toimitustodiste' -kortti.",
        "Katso noudossa ja toimituksessa otetut valokuvat.",
        "Napauta valokuvaa avataksesi sen koko n\u00E4yt\u00F6lle.",
    ], "Todistekuvat s\u00E4ilyv\u00E4t tilauksen yhteenvedossa kirjanpitoasi varten."),
    ("Laskut ja kuitit", [
        "Maksukuitit ja laskut luodaan automaattisesti PDF-tiedostoina.",
        "Ne l\u00E4hetet\u00E4\u00E4n s\u00E4hk\u00F6postiisi ja tallennetaan sovellukseen ladattavaksi milloin tahansa.",
    ], "Erinomaista kirjanpitoon \u2014 jokaisesta tapahtumasta on asiakirja."),
    ("Tallennetut maksutavat", [
        "Lis\u00E4\u00E4 ja hallitse kortteja kohdassa Asetukset \u2192 Maksutavat.",
        "Aseta oletuskortti ja poista vanhat.",
        "K\u00E4yt\u00E4 yhden napautuksen 'Maksa tallennetulla kortilla' nopeampaan maksuun.",
    ], None),
    ("Asetukset, kieli ja tuki", [
        "Hallitse yritysprofiiliasi ja asetuksiasi Asetuksissa.",
        "Vaihda sovelluksen kieli englannin ja suomen v\u00E4lill\u00E4 milloin tahansa.",
        "L\u00F6yd\u00E4t 'Ohje ja tuki' sek\u00E4 'Tietosuoja ja ehdot' profiilistasi.",
        "Tarvitsetko apua? Ota yhteytt\u00E4: care@nadaruns.com.",
    ], "Siirr\u00E4 enemm\u00E4n rahtia, v\u00E4hemm\u00E4ll\u00E4 tyhj\u00E4ll\u00E4 tilalla tiell\u00E4."),
]


if __name__ == "__main__":
    build(EN, "Driver User Guide", "How to earn with NadaRuns \u2014 from sign-up to payout", driver_en, "NadaRuns_Driver_Guide_EN.pptx")
    build(EN, "Shipper User Guide", "How to move freight with NadaRuns \u2014 from booking to delivery", shipper_en, "NadaRuns_Shipper_Guide_EN.pptx")
    build(FI, "Kuljettajan k\u00E4ytt\u00F6opas", "N\u00E4in ansaitset NadaRunsilla \u2014 rekister\u00F6innist\u00E4 maksuun", driver_fi, "NadaRuns_Kuljettajan_Opas_FI.pptx")
    build(FI, "L\u00E4hett\u00E4j\u00E4n k\u00E4ytt\u00F6opas", "N\u00E4in siirr\u00E4t rahtia NadaRunsilla \u2014 tilauksesta toimitukseen", shipper_fi, "NadaRuns_Lahettajan_Opas_FI.pptx")
    print("ALL DONE")
